"""
Outil Génération de documents — PDF, Word, Excel, PowerPoint.
"""
import os
from datetime import datetime

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "data", "outputs")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ─── PDF ────────────────────────────────────────────────────────────────────
def generate_pdf(title: str, content: str, filename: str = "") -> str:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.enums import TA_LEFT

        if not filename:
            filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        path = os.path.join(OUTPUT_DIR, filename)

        doc = SimpleDocTemplate(path, pagesize=A4,
                                leftMargin=2.5*cm, rightMargin=2.5*cm,
                                topMargin=2.5*cm, bottomMargin=2.5*cm)
        styles = getSampleStyleSheet()
        story = [
            Paragraph(title, styles["Title"]),
            Spacer(1, 0.5*cm),
        ]
        for para in content.split("\n\n"):
            if para.strip():
                story.append(Paragraph(para.replace("\n", "<br/>"), styles["Normal"]))
                story.append(Spacer(1, 0.3*cm))
        doc.build(story)
        return f"PDF généré : {path}"
    except ImportError:
        return "Erreur : installe reportlab (pip install reportlab)"
    except Exception as e:
        return f"Erreur PDF : {e}"


# ─── Word ────────────────────────────────────────────────────────────────────
def generate_word(title: str, content: str, filename: str = "") -> str:
    try:
        from docx import Document
        from docx.shared import Pt

        if not filename:
            filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        path = os.path.join(OUTPUT_DIR, filename)

        doc = Document()
        doc.add_heading(title, 0)
        for para in content.split("\n\n"):
            if para.strip():
                if para.startswith("## "):
                    doc.add_heading(para[3:], level=2)
                elif para.startswith("# "):
                    doc.add_heading(para[2:], level=1)
                else:
                    doc.add_paragraph(para)
        doc.save(path)
        return f"Word généré : {path}"
    except ImportError:
        return "Erreur : installe python-docx (pip install python-docx)"
    except Exception as e:
        return f"Erreur Word : {e}"


# ─── Excel ───────────────────────────────────────────────────────────────────
def generate_excel(title: str, data: list, headers: list = None, filename: str = "") -> str:
    """
    data : liste de listes (lignes)
    headers : liste de titres de colonnes
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        if not filename:
            filename = f"tableau_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        path = os.path.join(OUTPUT_DIR, filename)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31]

        header_fill = PatternFill("solid", fgColor="1E40AF")
        header_font = Font(bold=True, color="FFFFFF")

        start_row = 1
        if headers:
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
            start_row = 2

        for row_idx, row in enumerate(data, start_row):
            for col_idx, val in enumerate(row if isinstance(row, list) else [row], 1):
                ws.cell(row=row_idx, column=col_idx, value=val)

        # Auto-largeur colonnes
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=0)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)

        wb.save(path)
        return f"Excel généré : {path}"
    except ImportError:
        return "Erreur : installe openpyxl (pip install openpyxl)"
    except Exception as e:
        return f"Erreur Excel : {e}"


# ─── PowerPoint ──────────────────────────────────────────────────────────────
def generate_pptx(title: str, slides: list, filename: str = "") -> str:
    """
    slides : liste de dicts {"title": str, "content": str}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        if not filename:
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        path = os.path.join(OUTPUT_DIR, filename)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de titre
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = datetime.now().strftime("%d/%m/%Y")

        # Slides de contenu
        content_layout = prs.slide_layouts[1]
        for s in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = s.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            content = s.get("content", "")
            for i, line in enumerate(content.split("\n")):
                if line.strip():
                    if i == 0:
                        tf.text = line
                    else:
                        tf.add_paragraph().text = line

        prs.save(path)
        return f"PowerPoint généré : {path}"
    except ImportError:
        return "Erreur : installe python-pptx (pip install python-pptx)"
    except Exception as e:
        return f"Erreur PPTX : {e}"


# ─── Schémas ─────────────────────────────────────────────────────────────────
GENERATE_PDF_SCHEMA = {
    "type": "function",
    "function": {
        "name": "generate_pdf",
        "description": "Génère un fichier PDF avec un titre et du contenu texte.",
        "parameters": {
            "type": "object",
            "properties": {
                "title":    {"type": "string", "description": "Titre du document"},
                "content":  {"type": "string", "description": "Contenu du document (paragraphes séparés par \\n\\n)"},
                "filename": {"type": "string", "description": "Nom du fichier (optionnel, ex: rapport.pdf)"},
            },
            "required": ["title", "content"],
        },
    },
}

GENERATE_WORD_SCHEMA = {
    "type": "function",
    "function": {
        "name": "generate_word",
        "description": "Génère un fichier Word (.docx) avec un titre et du contenu. Supporte les titres Markdown (# et ##).",
        "parameters": {
            "type": "object",
            "properties": {
                "title":    {"type": "string", "description": "Titre du document"},
                "content":  {"type": "string", "description": "Contenu (supporte # Titre1, ## Titre2)"},
                "filename": {"type": "string", "description": "Nom du fichier (optionnel, ex: rapport.docx)"},
            },
            "required": ["title", "content"],
        },
    },
}

GENERATE_EXCEL_SCHEMA = {
    "type": "function",
    "function": {
        "name": "generate_excel",
        "description": "Génère un fichier Excel (.xlsx) avec des données tabulaires.",
        "parameters": {
            "type": "object",
            "properties": {
                "title":    {"type": "string", "description": "Nom de la feuille"},
                "data":     {"type": "array",  "description": "Données : liste de listes (lignes)", "items": {}},
                "headers":  {"type": "array",  "description": "Titres des colonnes (optionnel)", "items": {"type": "string"}},
                "filename": {"type": "string", "description": "Nom du fichier (optionnel, ex: tableau.xlsx)"},
            },
            "required": ["title", "data"],
        },
    },
}

GENERATE_PPTX_SCHEMA = {
    "type": "function",
    "function": {
        "name": "generate_pptx",
        "description": "Génère une présentation PowerPoint (.pptx).",
        "parameters": {
            "type": "object",
            "properties": {
                "title":    {"type": "string", "description": "Titre de la présentation"},
                "slides":   {"type": "array",  "description": "Liste de slides : [{title, content}]", "items": {"type": "object"}},
                "filename": {"type": "string", "description": "Nom du fichier (optionnel, ex: presentation.pptx)"},
            },
            "required": ["title", "slides"],
        },
    },
}
